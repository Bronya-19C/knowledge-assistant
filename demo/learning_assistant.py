# -*- coding: utf-8 -*-
import os
import json
import argparse
from docx import Document
from pptx import Presentation
from camel.agents import ChatAgent
from camel.configs import SiliconFlowConfig
from camel.models import ModelFactory
from camel.types import ModelPlatformType

class DocumentLoader:
    def load(self, path):
        ext = os.path.splitext(path)[1].lower()
        if ext == '.txt':
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        if ext in ('.doc', '.docx'):
            doc = Document(path)
            return '\n'.join(p.text for p in doc.paragraphs)
        if ext == '.pptx':
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        texts.append(shape.text)
            return '\n'.join(texts)
        raise ValueError(f'Unsupported format: {ext}')

class MemoryStore:
    def __init__(self, path='memory.json'):
        self.path = path
        if os.path.exists(path):
            with open(path, 'r', encoding='utf-8') as f:
                self.data = json.load(f)
            # 确保所有字段存在
            self.data.setdefault('domains', [])
            self.data.setdefault('documents', [])
        else:
            # 初始化包括 domains 的数据结构
            self.data = {'domains': [], 'documents': []}
            self._save()

    def add_domain(self, domain):
        if domain and domain not in self.data['domains']:
            self.data['domains'].append(domain)
            self._save()

    def add_document(self, text):
        if text and text not in self.data['documents']:
            self.data['documents'].append(text)
            self._save()

    def get_context(self):
        # 将 domains, documents 一并作为上下文
        domains = '\n'.join(self.data.get('domains', []))
        docs = '\n'.join(self.data.get('documents', []))
        return (
            f'已了解知识领域：\n{domains}\n'
            f'已提交文档：\n{docs}\n'
        )

    def _save(self):
        with open(self.path, 'w', encoding='utf-8') as f:
            json.dump(self.data, f, ensure_ascii=False, indent=2)

class LearningAgent:
    def __init__(self, api_key):
        # 主 Agent：生成面向学生的全面提纲
        sys_msg = (
            '你是一个学习助手，请用简洁的中文，根据学生身份，'
            '为学生生成面向学习的全面提纲，包括各章节或板块名称及对应主要知识点摘要，'
            '不要输出教学目标、课时安排或其他教学设计内容。'
        )
        model = ModelFactory.create(
            model_platform=ModelPlatformType.SILICONFLOW,
            model_type='deepseek-ai/DeepSeek-V3',
            model_config_dict=SiliconFlowConfig(stream=True, temperature=0.3).as_dict(),
            api_key=api_key
        )
        self.agent = ChatAgent(system_message=sys_msg, model=model, output_language='中文')

    def extract_domains(self, content, memory_context):
        # 使用主 Agent 提取该文档的知识领域
        prompt = (
            f"{memory_context}\n"
            f"请从以下内容中提取简洁的知识领域关键词列表（以逗号或换行分隔）：\n{content}"
        )
        resp = self.agent.step(prompt).msg.content
        # 粗略拆分关键词
        domains = [d.strip() for d in resp.replace('，', ',').split(',') if d.strip()]
        return domains

    def generate_outline(self, content, memory_context):
        # 生成学生全面学习提纲
        prompt = (
            f"{memory_context}\n"
            f"请根据以下内容，为学生生成全面的学习提纲，列出每个章节或板块标题及其主要知识点：\n{content}"
        )
        return self.agent.step(prompt).msg.content

if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        description='加载文档并为学生生成全面学习提纲，并记忆知识领域'
    )
    parser.add_argument('files', nargs='+', help='输入文件路径，支持 txt/docx/pptx')
    parser.add_argument('--api_key', required=True, help='SiliconFlow API 密钥')
    args = parser.parse_args()

    loader = DocumentLoader()
    mem = MemoryStore()
    agent = LearningAgent(api_key=args.api_key)

    # 逐个加载文档，记忆知识领域
    for path in args.files:
        print(f"加载文件: {path}")
        text = loader.load(path)
        mem.add_document(text)
        # 提取并记忆该文档的知识领域
        domains = agent.extract_domains(text, mem.get_context())
        for d in domains:
            mem.add_domain(d)

    # 使用最后一个文档生成提纲
    content = mem.data['documents'][-1]
    outline = agent.generate_outline(content, mem.get_context())
    print("学生学习提纲:")
    print(outline)
